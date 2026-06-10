#!/usr/bin/python3
"""Assignment 4: Tree and ensemble classifiers on three datasets."""

import json
import subprocess
import textwrap
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import requests

# python-pptx 0.6.x can trip a Python 3.6 textwrap edge case while generating
# enum documentation strings at import time. The patch is local to this process.
_ORIGINAL_TEXTWRAP_FILL = textwrap.fill


def _safe_textwrap_fill(text, *args, **kwargs):
    width = kwargs.get("width", 70)
    subsequent_indent = kwargs.get("subsequent_indent", "")
    if width <= len(subsequent_indent):
        kwargs["width"] = len(subsequent_indent) + 20
    return _ORIGINAL_TEXTWRAP_FILL(text, *args, **kwargs)


textwrap.fill = _safe_textwrap_fill
from pptx import Presentation
from pptx.util import Inches, Pt
from scipy import sparse
from sklearn.base import clone
from sklearn.compose import ColumnTransformer
from sklearn.decomposition import PCA
from sklearn.ensemble import AdaBoostClassifier, RandomForestClassifier
from sklearn.impute import SimpleImputer
from sklearn.metrics import (
    accuracy_score,
    auc,
    f1_score,
    precision_score,
    recall_score,
    roc_auc_score,
    roc_curve,
)
from sklearn.model_selection import StratifiedKFold, train_test_split
from sklearn.pipeline import Pipeline
from sklearn.preprocessing import LabelEncoder, OneHotEncoder, StandardScaler, label_binarize
from sklearn.tree import DecisionTreeClassifier

try:
    from xgboost import XGBClassifier

    XGBOOST_AVAILABLE = True
except Exception:
    from sklearn.ensemble import GradientBoostingClassifier

    XGBOOST_AVAILABLE = False


BASE_DIR = Path(__file__).resolve().parent
DATA_DIR = BASE_DIR / "data"
IMAGE_DIR = BASE_DIR / "images"
RESULT_DIR = BASE_DIR / "results"
RANDOM_STATE = 42

DATASETS = {
    "cervical": {
        "name": "Cervical Cancer Risk Factors",
        "url": "https://archive.ics.uci.edu/ml/machine-learning-databases/00383/risk_factors_cervical_cancer.csv",
        "file": "risk_factors_cervical_cancer.csv",
        "target": "Biopsy",
        "source": "https://christophm.github.io/interpretable-ml-book/cervical.html",
    },
    "fetal_health": {
        "name": "Fetal Health Classification",
        "url": "https://raw.githubusercontent.com/SagarSharma4244/Fetal-Health/main/fetal_health.csv",
        "file": "fetal_health.csv",
        "target": "fetal_health",
        "source": "https://www.kaggle.com/datasets/andrewmvd/fetal-health-classification",
    },
    "banking": {
        "name": "Banking Marketing",
        "url": "https://raw.githubusercontent.com/AndrzejSzymanski/TDS/master/banking.csv",
        "file": "banking.csv",
        "target": None,
        "source": "https://github.com/AndrzejSzymanski/TDS/blob/master/banking.csv",
    },
}


def download(url, path):
    DATA_DIR.mkdir(exist_ok=True)
    if path.exists():
        return path
    response = requests.get(url, timeout=120)
    response.raise_for_status()
    path.write_bytes(response.content)
    return path


def cap_iqr_outliers(X):
    X = X.copy()
    numeric_cols = X.select_dtypes(include=[np.number]).columns
    outlier_counts = {}
    for col in numeric_cols:
        series = X[col]
        q1 = series.quantile(0.25)
        q3 = series.quantile(0.75)
        iqr = q3 - q1
        if pd.isna(iqr) or iqr == 0:
            outlier_counts[col] = 0
            continue
        lower = q1 - 1.5 * iqr
        upper = q3 + 1.5 * iqr
        outlier_counts[col] = int(((series < lower) | (series > upper)).sum())
        X[col] = series.clip(lower=lower, upper=upper)
    return X, outlier_counts


def load_cervical():
    spec = DATASETS["cervical"]
    path = download(spec["url"], DATA_DIR / spec["file"])
    df = pd.read_csv(path, na_values="?")
    y = df["Biopsy"].astype(int).map({0: "Healthy", 1: "Cancer"})
    leakage_cols = ["Biopsy", "Hinselmann", "Schiller", "Citology", "Dx:Cancer", "Dx:CIN", "Dx:HPV", "Dx"]
    X = df.drop(columns=[col for col in leakage_cols if col in df.columns])
    for col in X.columns:
        X[col] = pd.to_numeric(X[col], errors="coerce")
    X, outliers = cap_iqr_outliers(X)
    return {
        "key": "cervical",
        "name": spec["name"],
        "source": spec["source"],
        "X": X,
        "y": y,
        "outliers": outliers,
        "notes": "Missing values marked '?' were median-imputed. Diagnostic target/leakage columns were removed from features.",
    }


def load_fetal_health():
    spec = DATASETS["fetal_health"]
    path = download(spec["url"], DATA_DIR / spec["file"])
    df = pd.read_csv(path)
    mapping = {1.0: "Normal", 2.0: "Suspect", 3.0: "Pathological"}
    y = df[spec["target"]].map(mapping)
    X = df.drop(columns=[spec["target"]])
    X, outliers = cap_iqr_outliers(X)
    return {
        "key": "fetal_health",
        "name": spec["name"],
        "source": spec["source"],
        "X": X,
        "y": y,
        "outliers": outliers,
        "notes": "Numerical CTG features were IQR-clipped and median-imputed.",
    }


def load_banking():
    spec = DATASETS["banking"]
    path = download(spec["url"], DATA_DIR / spec["file"])
    df = pd.read_csv(path)
    target = df.columns[-1]
    y = df[target].map({0: "No", 1: "Yes"}).fillna(df[target].astype(str))
    X = df.drop(columns=[target])
    numeric_cols = X.select_dtypes(include=[np.number]).columns
    for col in X.columns:
        if col not in numeric_cols:
            X[col] = X[col].astype(str).replace("unknown", np.nan)
    X, outliers = cap_iqr_outliers(X)
    return {
        "key": "banking",
        "name": spec["name"],
        "source": spec["source"],
        "X": X,
        "y": y,
        "outliers": outliers,
        "notes": "Categorical 'unknown' values were treated as missing and imputed before one-hot encoding.",
    }


def load_all_datasets():
    return [load_cervical(), load_fetal_health(), load_banking()]


def make_preprocessor(X):
    numeric_cols = list(X.select_dtypes(include=[np.number]).columns)
    categorical_cols = [col for col in X.columns if col not in numeric_cols]
    transformers = []
    if numeric_cols:
        transformers.append(
            (
                "num",
                Pipeline(
                    [
                        ("imputer", SimpleImputer(strategy="median")),
                        ("scaler", StandardScaler()),
                    ]
                ),
                numeric_cols,
            )
        )
    if categorical_cols:
        transformers.append(
            (
                "cat",
                Pipeline(
                    [
                        ("imputer", SimpleImputer(strategy="most_frequent")),
                        ("onehot", OneHotEncoder(handle_unknown="ignore")),
                    ]
                ),
                categorical_cols,
            )
        )
    return ColumnTransformer(transformers)


def get_model(model_name, n_classes):
    if model_name == "Decision Tree":
        return DecisionTreeClassifier(
            random_state=RANDOM_STATE, class_weight="balanced", min_samples_leaf=3
        ), "DecisionTreeClassifier"
    if model_name == "Random Forest":
        return RandomForestClassifier(
            n_estimators=80,
            random_state=RANDOM_STATE,
            class_weight="balanced",
            n_jobs=-1,
            min_samples_leaf=2,
        ), "RandomForestClassifier"
    if model_name == "XGBoost":
        if XGBOOST_AVAILABLE:
            objective = "multi:softprob" if n_classes > 2 else "binary:logistic"
            eval_metric = "mlogloss" if n_classes > 2 else "logloss"
            return XGBClassifier(
                n_estimators=80,
                max_depth=3,
                learning_rate=0.08,
                subsample=0.9,
                colsample_bytree=0.9,
                objective=objective,
                eval_metric=eval_metric,
                random_state=RANDOM_STATE,
                n_jobs=1,
                use_label_encoder=False,
                verbosity=0,
            ), "XGBClassifier"
        return GradientBoostingClassifier(random_state=RANDOM_STATE), "GradientBoosting fallback"
    if model_name == "AdaBoost":
        base = DecisionTreeClassifier(max_depth=2, random_state=RANDOM_STATE)
        return AdaBoostClassifier(
            base_estimator=base,
            n_estimators=80,
            learning_rate=0.5,
            random_state=RANDOM_STATE,
        ), "AdaBoostClassifier"
    raise ValueError("Unknown model: {}".format(model_name))


def calculate_metrics(y_true, y_pred, y_proba, n_classes):
    metrics = {
        "precision_weighted": float(precision_score(y_true, y_pred, average="weighted", zero_division=0)),
        "recall_weighted": float(recall_score(y_true, y_pred, average="weighted", zero_division=0)),
        "f1_weighted": float(f1_score(y_true, y_pred, average="weighted", zero_division=0)),
        "accuracy": float(accuracy_score(y_true, y_pred)),
        "auc_roc": np.nan,
    }
    try:
        if n_classes == 2:
            metrics["auc_roc"] = float(roc_auc_score(y_true, y_proba[:, 1]))
        else:
            metrics["auc_roc"] = float(
                roc_auc_score(y_true, y_proba, average="weighted", multi_class="ovr")
            )
    except Exception:
        metrics["auc_roc"] = np.nan
    return metrics


def evaluate_dataset(dataset):
    X = dataset["X"]
    y_labels = dataset["y"]
    encoder = LabelEncoder()
    y = encoder.fit_transform(y_labels)
    n_classes = len(encoder.classes_)
    X_train, X_test, y_train, y_test = train_test_split(
        X, y, test_size=0.20, random_state=RANDOM_STATE, stratify=y
    )
    class_distribution = {
        str(cls): int(count)
        for cls, count in zip(*np.unique(y_labels, return_counts=True))
    }

    rows = []
    roc_payload = {}
    fitted_payload = {}
    cv = StratifiedKFold(n_splits=5, shuffle=True, random_state=RANDOM_STATE)

    for model_name in ["Decision Tree", "Random Forest", "XGBoost", "AdaBoost"]:
        fold_metrics = []
        backend = None
        for fold, (fit_idx, val_idx) in enumerate(cv.split(X_train, y_train), start=1):
            model, backend = get_model(model_name, n_classes)
            pipe = Pipeline(
                [("preprocess", make_preprocessor(X_train)), ("model", model)]
            )
            pipe.fit(X_train.iloc[fit_idx], y_train[fit_idx])
            pred = pipe.predict(X_train.iloc[val_idx])
            proba = pipe.predict_proba(X_train.iloc[val_idx])
            metrics = calculate_metrics(y_train[val_idx], pred, proba, n_classes)
            metrics.update(
                {
                    "dataset": dataset["key"],
                    "model": model_name,
                    "backend": backend,
                    "split": "cv",
                    "fold": fold,
                }
            )
            fold_metrics.append(metrics)
            rows.append(metrics)

        model, backend = get_model(model_name, n_classes)
        final_pipe = Pipeline(
            [("preprocess", make_preprocessor(X_train)), ("model", model)]
        )
        final_pipe.fit(X_train, y_train)
        test_pred = final_pipe.predict(X_test)
        test_proba = final_pipe.predict_proba(X_test)
        test_metrics = calculate_metrics(y_test, test_pred, test_proba, n_classes)
        test_metrics.update(
            {
                "dataset": dataset["key"],
                "model": model_name,
                "backend": backend,
                "split": "test",
                "fold": 0,
            }
        )
        rows.append(test_metrics)
        roc_payload[model_name] = {
            "y_test": y_test,
            "proba": test_proba,
            "classes": list(encoder.classes_),
            "n_classes": n_classes,
        }
        fitted_payload[model_name] = {
            "model": final_pipe,
            "X_train": X_train,
            "X_test": X_test,
            "y_train": y_train,
            "y_test": y_test,
            "classes": list(encoder.classes_),
            "n_classes": n_classes,
        }
        print(
            "{} / {} test accuracy={:.4f} AUC={:.4f}".format(
                dataset["key"], model_name, test_metrics["accuracy"], test_metrics["auc_roc"]
            )
        )

    profile = {
        "dataset": dataset["key"],
        "name": dataset["name"],
        "source": dataset["source"],
        "rows": int(len(X)),
        "features": int(X.shape[1]),
        "classes": list(encoder.classes_),
        "class_distribution": class_distribution,
        "train_rows": int(len(X_train)),
        "test_rows": int(len(X_test)),
        "missing_values": {col: int(X[col].isna().sum()) for col in X.columns},
        "iqr_outlier_counts": dataset["outliers"],
        "notes": dataset["notes"],
    }
    return pd.DataFrame(rows), profile, roc_payload, fitted_payload


def plot_roc_curves(all_roc):
    IMAGE_DIR.mkdir(exist_ok=True)
    datasets = list(all_roc.keys())
    models = ["Decision Tree", "Random Forest", "XGBoost", "AdaBoost"]
    fig, axes = plt.subplots(len(datasets), len(models), figsize=(18, 12))
    for row_idx, dataset_key in enumerate(datasets):
        for col_idx, model_name in enumerate(models):
            ax = axes[row_idx, col_idx]
            payload = all_roc[dataset_key][model_name]
            y_test = payload["y_test"]
            proba = payload["proba"]
            n_classes = payload["n_classes"]
            if n_classes == 2:
                fpr, tpr, _ = roc_curve(y_test, proba[:, 1])
                score = auc(fpr, tpr)
            else:
                y_bin = label_binarize(y_test, classes=list(range(n_classes)))
                fpr, tpr, _ = roc_curve(y_bin.ravel(), proba.ravel())
                score = auc(fpr, tpr)
            ax.plot(fpr, tpr, color="#4C78A8", linewidth=2, label="AUC={:.3f}".format(score))
            ax.plot([0, 1], [0, 1], color="#BAB0AC", linestyle="--", linewidth=1)
            ax.set_title("{} / {}".format(dataset_key, model_name), fontsize=9)
            ax.set_xlabel("FPR")
            ax.set_ylabel("TPR")
            ax.legend(fontsize=8)
    fig.tight_layout()
    fig.savefig(IMAGE_DIR / "roc_curves.png", dpi=160)
    plt.close(fig)


def as_dense(matrix):
    if sparse.issparse(matrix):
        return matrix.toarray()
    return np.asarray(matrix)


def plot_decision_boundaries(all_fitted):
    datasets = list(all_fitted.keys())
    models = ["Decision Tree", "Random Forest", "XGBoost", "AdaBoost"]
    fig, axes = plt.subplots(len(datasets), len(models), figsize=(18, 12))
    for row_idx, dataset_key in enumerate(datasets):
        for col_idx, model_name in enumerate(models):
            ax = axes[row_idx, col_idx]
            payload = all_fitted[dataset_key][model_name]
            pipe = payload["model"]
            X_train = payload["X_train"]
            X_test = payload["X_test"]
            y_train = payload["y_train"]
            y_test = payload["y_test"]
            n_classes = payload["n_classes"]

            pre = clone(pipe.named_steps["preprocess"])
            X_train_pre = as_dense(pre.fit_transform(X_train))
            X_test_pre = as_dense(pre.transform(X_test))
            pca = PCA(n_components=2, random_state=RANDOM_STATE)
            X_train_2d = pca.fit_transform(X_train_pre)
            X_test_2d = pca.transform(X_test_pre)

            model, _ = get_model(model_name, n_classes)
            model.fit(X_train_2d, y_train)
            x_min, x_max = X_train_2d[:, 0].min() - 1, X_train_2d[:, 0].max() + 1
            y_min, y_max = X_train_2d[:, 1].min() - 1, X_train_2d[:, 1].max() + 1
            xx, yy = np.meshgrid(
                np.linspace(x_min, x_max, 140), np.linspace(y_min, y_max, 140)
            )
            Z = model.predict(np.c_[xx.ravel(), yy.ravel()]).reshape(xx.shape)
            ax.contourf(xx, yy, Z, alpha=0.22, cmap="tab10")
            rng = np.random.RandomState(RANDOM_STATE)
            train_sample = rng.choice(
                len(X_train_2d), size=min(600, len(X_train_2d)), replace=False
            )
            test_sample = rng.choice(
                len(X_test_2d), size=min(250, len(X_test_2d)), replace=False
            )
            ax.scatter(
                X_train_2d[train_sample, 0],
                X_train_2d[train_sample, 1],
                c=y_train[train_sample],
                s=8,
                cmap="tab10",
                alpha=0.35,
                edgecolors="none",
            )
            ax.scatter(
                X_test_2d[test_sample, 0],
                X_test_2d[test_sample, 1],
                c=y_test[test_sample],
                s=14,
                cmap="tab10",
                alpha=0.75,
                marker="x",
            )
            ax.set_title("{} / {}".format(dataset_key, model_name), fontsize=9)
            ax.set_xlabel("PC1")
            ax.set_ylabel("PC2")
    fig.tight_layout()
    fig.savefig(IMAGE_DIR / "decision_boundaries.png", dpi=160)
    plt.close(fig)


def plot_metric_summary(metrics):
    test = metrics[metrics["split"] == "test"].copy()
    fig, ax = plt.subplots(figsize=(11, 5))
    test["label"] = test["dataset"] + "\n" + test["model"]
    ax.bar(test["label"], test["accuracy"], color="#4C78A8")
    ax.set_ylim(0, 1.05)
    ax.set_ylabel("Test accuracy")
    ax.set_title("Test Accuracy by Dataset and Classifier")
    ax.tick_params(axis="x", rotation=60)
    fig.tight_layout()
    fig.savefig(IMAGE_DIR / "test_accuracy_summary.png", dpi=160)
    plt.close(fig)


def latex_table(df, columns):
    return df[columns].to_latex(index=False, escape=True, float_format=lambda x: "{:.4f}".format(x))


def write_readme_pdf(metrics, profiles):
    test = metrics[metrics["split"] == "test"].copy()
    cv = (
        metrics[metrics["split"] == "cv"]
        .groupby(["dataset", "model", "backend"], as_index=False)[
            ["precision_weighted", "recall_weighted", "f1_weighted", "accuracy", "auc_roc"]
        ]
        .mean()
    )
    readme = BASE_DIR / "readme.md"
    readme.write_text(
        "# Assignment 4 Readme\n\n"
        "Run with `/usr/bin/python3 assignment4_classification.py`.\n\n"
        "Outputs are written to `results/` and `images/`. The implementation uses stratified 80/20 train/test splits and 5-fold cross-validation on the training split.\n"
    )

    profile_rows = [
        {
            "dataset": p["dataset"],
            "rows": p["rows"],
            "features": p["features"],
            "train_rows": p["train_rows"],
            "test_rows": p["test_rows"],
            "classes": ", ".join(map(str, p["classes"])),
        }
        for p in profiles
    ]
    profile_df = pd.DataFrame(profile_rows)

    tex = [
        r"\documentclass[11pt]{article}",
        r"\usepackage[margin=0.7in]{geometry}",
        r"\usepackage{graphicx}",
        r"\usepackage{booktabs}",
        r"\usepackage{float}",
        r"\begin{document}",
        r"\title{Data Mining Assignment 4: Classification Readme}",
        r"\author{}",
        r"\date{}",
        r"\maketitle",
        r"\section*{Datasets and Preprocessing}",
        "All datasets were split in an 80/20 stratified ratio. Missing values were imputed inside the training pipeline. Numeric outliers were clipped using the 1.5 IQR rule before modeling.",
        profile_df.to_latex(index=False, escape=True),
        r"\section*{5-Fold Cross-Validation Mean Metrics}",
        latex_table(cv, ["dataset", "model", "backend", "precision_weighted", "recall_weighted", "f1_weighted", "accuracy", "auc_roc"]),
        r"\section*{Held-Out Test Metrics}",
        latex_table(test, ["dataset", "model", "backend", "precision_weighted", "recall_weighted", "f1_weighted", "accuracy", "auc_roc"]),
        r"\section*{Figures}",
        r"\begin{figure}[H]\centering\includegraphics[width=\textwidth]{images/test_accuracy_summary.png}\caption{Test accuracy summary.}\end{figure}",
        r"\begin{figure}[H]\centering\includegraphics[width=\textwidth]{images/roc_curves.png}\caption{ROC curves.}\end{figure}",
        r"\begin{figure}[H]\centering\includegraphics[width=\textwidth]{images/decision_boundaries.png}\caption{PCA decision-boundary visualizations.}\end{figure}",
        r"\end{document}",
    ]
    (BASE_DIR / "readme.tex").write_text("\n".join(tex))
    try:
        subprocess.run(
            ["pdflatex", "-interaction=nonstopmode", "readme.tex"],
            cwd=str(BASE_DIR),
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
        )
    except Exception as exc:
        print("Warning: readme PDF generation failed: {}".format(exc))


def add_title(slide, title, subtitle=None):
    slide.shapes.title.text = title
    if subtitle:
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.8), Inches(1.0))
        tx.text_frame.text = subtitle


def write_ppt(metrics):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[5]

    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Assignment 4 Classification"
    slide.placeholders[1].text = "Decision Tree, Random Forest, XGBoost, and AdaBoost"

    slide = prs.slides.add_slide(content_layout)
    add_title(slide, "Workflow", "Data cleaning -> stratified 80/20 split -> 5-fold CV -> held-out test evaluation.")
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(8.8), Inches(2.4))
    frame = box.text_frame
    frame.text = "Preprocessing"
    for item in [
        "Median/mode imputation inside the modeling pipeline",
        "IQR clipping for numeric outliers",
        "One-hot encoding for categorical banking features",
        "PCA projection used only for visual decision-boundary plots",
    ]:
        p = frame.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(18)

    for title, image_name in [
        ("Test Accuracy", "test_accuracy_summary.png"),
        ("ROC Curves", "roc_curves.png"),
        ("Decision Boundaries", "decision_boundaries.png"),
    ]:
        slide = prs.slides.add_slide(content_layout)
        add_title(slide, title)
        slide.shapes.add_picture(str(IMAGE_DIR / image_name), Inches(0.4), Inches(1.2), width=Inches(9.2))

    test = metrics[metrics["split"] == "test"].sort_values(["dataset", "accuracy"], ascending=[True, False])
    best = test.groupby("dataset", as_index=False).first()
    slide = prs.slides.add_slide(content_layout)
    add_title(slide, "Best Test Models")
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.8), Inches(4.5))
    frame = box.text_frame
    frame.text = "Top held-out results"
    for _, row in best.iterrows():
        p = frame.add_paragraph()
        p.text = "{}: {} accuracy {:.3f}, AUC {:.3f}".format(
            row["dataset"], row["model"], row["accuracy"], row["auc_roc"]
        )
        p.level = 1
        p.font.size = Pt(18)

    prs.save(BASE_DIR / "Assignment4_classification.pptx")


def save_outputs(metrics, profiles):
    RESULT_DIR.mkdir(exist_ok=True)
    metrics.to_csv(RESULT_DIR / "classification_metrics.csv", index=False)
    with open(RESULT_DIR / "dataset_profiles.json", "w") as fh:
        json.dump(profiles, fh, indent=2, sort_keys=True)
    with open(RESULT_DIR / "assignment4_results.json", "w") as fh:
        json.dump(
            {
                "xgboost_available": XGBOOST_AVAILABLE,
                "datasets": profiles,
                "test_metrics": metrics[metrics["split"] == "test"].to_dict(orient="records"),
            },
            fh,
            indent=2,
            sort_keys=True,
        )


def main():
    datasets = load_all_datasets()
    all_metrics = []
    profiles = []
    all_roc = {}
    all_fitted = {}
    for dataset in datasets:
        metrics, profile, roc_payload, fitted_payload = evaluate_dataset(dataset)
        all_metrics.append(metrics)
        profiles.append(profile)
        all_roc[dataset["key"]] = roc_payload
        all_fitted[dataset["key"]] = fitted_payload

    metrics = pd.concat(all_metrics, ignore_index=True)
    plot_roc_curves(all_roc)
    plot_decision_boundaries(all_fitted)
    plot_metric_summary(metrics)
    save_outputs(metrics, profiles)
    write_readme_pdf(metrics, profiles)
    write_ppt(metrics)
    print("Assignment 4 completed.")
    print("XGBoost available: {}".format(XGBOOST_AVAILABLE))
    print("Metrics: {}".format(RESULT_DIR / "classification_metrics.csv"))
    print("PPT: {}".format(BASE_DIR / "Assignment4_classification.pptx"))


if __name__ == "__main__":
    main()
